import asyncio
import html
import io
import logging
import os
import re
import sqlite3
import uuid
import json
import requests
import tempfile
import zipfile
from datetime import date
from telegram import InlineKeyboardButton, InlineKeyboardMarkup, Update
from telegram.ext import (
    ApplicationBuilder,
    CallbackQueryHandler,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    filters,
)

# ========= НАСТРОЙКИ =========
BOT_TOKEN = os.environ["BOT_TOKEN"]
AI_TOKEN = os.environ["GROK_TOKEN"]
CRYPTO_TOKEN = os.environ["CRYPTO_TOKEN"]
OPENROUTER_TOKEN = os.environ.get("OPENROUTER_TOKEN", "")
IMAGE_API_KEY = os.environ.get("IMAGE_API_KEY", "")
FISH_AUDIO_TOKEN = os.environ.get("FISH_AUDIO_TOKEN", "")

ADMINS = [8166720202, 1881900547, 8294681123]

CRYPTO_API = "https://pay.crypt.bot/api/"
WELCOME_IMAGE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "attached_assets", "089e36149455eb432c2afd94fe5f4bd8_1775834811270.jpg")
if not os.path.exists(WELCOME_IMAGE):
    WELCOME_IMAGE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "welcome.jpg")

# ========= БД =========
_DB_PATH = os.environ.get("DB_PATH", os.path.join(os.path.dirname(__file__), "bot.db"))
conn = sqlite3.connect(_DB_PATH, check_same_thread=False)
cursor = conn.cursor()

cursor.execute("""
CREATE TABLE IF NOT EXISTS users (
    user_id INTEGER PRIMARY KEY,
    requests INTEGER DEFAULT 10,
    referrals INTEGER DEFAULT 0,
    referrer INTEGER,
    banned INTEGER DEFAULT 0,
    model TEXT DEFAULT 'llama-3.3-70b-versatile',
    last_daily TEXT DEFAULT '',
    total_used INTEGER DEFAULT 0
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS payments (
    invoice_id TEXT,
    user_id INTEGER,
    amount INTEGER,
    status TEXT
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS promo_codes (
    code TEXT PRIMARY KEY,
    requests INTEGER,
    max_uses INTEGER,
    uses_count INTEGER DEFAULT 0
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS promo_uses (
    code TEXT,
    user_id INTEGER,
    PRIMARY KEY (code, user_id)
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS projects (
    project_id TEXT PRIMARY KEY,
    user_id INTEGER,
    title TEXT,
    description TEXT,
    files_json TEXT,
    created_at TEXT DEFAULT ''
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS user_memory (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    fact TEXT,
    created_at TEXT DEFAULT ''
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS reminders (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    remind_at TEXT,
    text TEXT,
    sent INTEGER DEFAULT 0
)
""")

for col_sql in [
    "ALTER TABLE users ADD COLUMN model TEXT DEFAULT 'llama-3.3-70b-versatile'",
    "ALTER TABLE users ADD COLUMN last_daily TEXT DEFAULT ''",
    "ALTER TABLE users ADD COLUMN total_used INTEGER DEFAULT 0",
]:
    try:
        cursor.execute(col_sql)
        conn.commit()
    except Exception:
        pass

conn.commit()

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# ========= ДОСТУПНЫЕ МОДЕЛИ =========
GROQ_MODELS = {
    "llama-3.3-70b-versatile": "🏴 Llama 3.3 70B [Groq] — умный",
    "llama-3.1-8b-instant":    "🏴 Llama 3.1 8B [Groq] — быстрый",
    "mixtral-8x7b-32768":      "🏴 Mixtral 8x7B [Groq] — длинный контекст",
    "gemma2-9b-it":            "🏴 Gemma 2 9B [Groq] — от Google",
}

OPENROUTER_MODELS = {
    "deepseek/deepseek-r1:free":              "🏴 DeepSeek R1 — мощное мышление",
    "deepseek/deepseek-chat-v3-0324:free":    "🏴 DeepSeek V3 — умный чат",
    "meta-llama/llama-4-scout:free":          "🏴 Llama 4 Scout — быстрый",
    "meta-llama/llama-4-maverick:free":       "🏴 Llama 4 Maverick — умный",
    "google/gemma-3-27b-it:free":             "🏴 Gemma 3 27B — от Google",
    "qwen/qwen3-235b-a22b:free":              "🏴 Qwen3 235B — огромный",
    "qwen/qwen3-30b-a3b:free":               "🏴 Qwen3 30B — баланс",
    "microsoft/phi-4:free":                   "🏴 Phi-4 [Microsoft] — компактный",
    "mistralai/mistral-7b-instruct:free":     "🏴 Mistral 7B — классика",
    "nvidia/llama-3.1-nemotron-70b-instruct:free": "🏴 Nemotron 70B [NVIDIA]",
}

MODELS = {**GROQ_MODELS, **OPENROUTER_MODELS}
DEFAULT_MODEL = "llama-3.3-70b-versatile"

def is_openrouter_model(model_id: str) -> bool:
    return "/" in model_id

# ========= КЛЮЧЕВЫЕ СЛОВА =========
ZIP_KEYWORDS = [
    "zip", "зип", "архив", "скачать", "скачай", "создай файл", "напиши файл",
    "сделай файл", "пришли файл", "отправь файл", "в файле", "файлом",
    "создай скрипт", "напиши скрипт", "сделай скрипт", "пришли скрипт",
    "игру", "игра", "game", "напиши игр", "сделай игр", "создай игр",
    "проект", "project", "приложение", "app", "сайт", "website",
    "напиши код", "сделай код", "создай код",
    "godot", "ren'py", "renpy", "unity", "pygame", "love2d", "love ",
    "движок", "визуальная новелла", "visual novel", "платформер", "platformer",
    "шутер", "shooter", "rpg", "рпг", "аркада", "arcade",
]

IMAGE_KEYWORDS = [
    "нарисуй", "нарисуй мне", "сгенерируй картинку", "сгенерируй изображение",
    "создай картинку", "создай изображение", "generate image", "draw", "draw me",
    "нарисуй картинку", "создай фото", "сделай картинку", "сделай изображение",
    "картинка", "изображение по запросу", "арт", "generate art", "create image",
    "create picture", "рисунок", "покажи картинку",
]

WORD_KEYWORDS = [
    "word", "docx", "документ word", "создай документ", "сделай документ",
    "напиши документ", "word файл", "doc файл",
]

PPT_KEYWORDS = [
    "powerpoint", "pptx", "презентация", "создай презентацию", "сделай презентацию",
    "напиши презентацию", "слайды", "slides", "ppt файл",
]

SURVEY_KEYWORDS = [
    "опрос", "викторина", "quiz", "questionnaire", "тест", "вопросы и ответы",
    "создай опрос", "сделай опрос", "создай тест", "сделай тест",
]

REMEMBER_KEYWORDS = [
    "запомни", "вспомни", "загрузи проект", "открой проект", "продолжи проект",
    "remember", "load project", "continue project",
]

EXTENSIONS = {
    "python": "py", "py": "py",
    "javascript": "js", "js": "js",
    "typescript": "ts", "ts": "ts",
    "html": "html", "htm": "html",
    "css": "css",
    "json": "json",
    "yaml": "yaml", "yml": "yaml",
    "bash": "sh", "sh": "sh", "shell": "sh",
    "sql": "sql",
    "php": "php",
    "java": "java",
    "cpp": "cpp", "c++": "cpp",
    "c": "c",
    "rust": "rs",
    "go": "go",
    "ruby": "rb", "rb": "rb",
    "kotlin": "kt",
    "swift": "swift",
    "lua": "lua",
    "xml": "xml",
    "toml": "toml",
    "ini": "ini",
    "markdown": "md", "md": "md",
    "plaintext": "txt", "text": "txt",
    "csharp": "cs", "cs": "cs", "c#": "cs",
    "gdscript": "gd", "gd": "gd",
    "tscn": "tscn",
    "tres": "tres",
    "renpy": "rpy", "rpy": "rpy",
    "cfg": "cfg",
}

def wants_zip(text: str) -> bool:
    t = text.lower()
    return any(kw in t for kw in ZIP_KEYWORDS)

def wants_image(text: str) -> bool:
    t = text.lower()
    return any(kw in t for kw in IMAGE_KEYWORDS)

def wants_word(text: str) -> bool:
    t = text.lower()
    return any(kw in t for kw in WORD_KEYWORDS)

def wants_ppt(text: str) -> bool:
    t = text.lower()
    return any(kw in t for kw in PPT_KEYWORDS)

def wants_survey(text: str) -> bool:
    t = text.lower()
    return any(kw in t for kw in SURVEY_KEYWORDS)

def wants_remember(text: str) -> bool:
    t = text.lower()
    return any(kw in t for kw in REMEMBER_KEYWORDS)

# ========= ГЕНЕРАЦИЯ PROJECT ID =========
def generate_project_id() -> str:
    return str(uuid.uuid4())[:8].upper()

def add_project_id_to_code(code: str, project_id: str, filename: str) -> str:
    """Добавляет ID проекта в конец файла в виде комментария."""
    ext = filename.rsplit(".", 1)[1].lower() if "." in filename else ""
    comment_styles = {
        "py": f"\n\n# ID: {project_id}",
        "gd": f"\n\n# ID: {project_id}",
        "js": f"\n\n// ID: {project_id}",
        "ts": f"\n\n// ID: {project_id}",
        "cs": f"\n\n// ID: {project_id}",
        "java": f"\n\n// ID: {project_id}",
        "cpp": f"\n\n// ID: {project_id}",
        "c": f"\n\n// ID: {project_id}",
        "go": f"\n\n// ID: {project_id}",
        "rs": f"\n\n// ID: {project_id}",
        "php": f"\n\n// ID: {project_id}",
        "kt": f"\n\n// ID: {project_id}",
        "swift": f"\n\n// ID: {project_id}",
        "lua": f"\n\n-- ID: {project_id}",
        "rb": f"\n\n# ID: {project_id}",
        "sh": f"\n\n# ID: {project_id}",
        "sql": f"\n\n-- ID: {project_id}",
        "html": f"\n\n<!-- ID: {project_id} -->",
        "xml": f"\n\n<!-- ID: {project_id} -->",
        "tscn": f"\n\n; ID: {project_id}",
        "tres": f"\n\n; ID: {project_id}",
        "cfg": f"\n\n; ID: {project_id}",
        "ini": f"\n\n; ID: {project_id}",
        "toml": f"\n\n# ID: {project_id}",
        "yaml": f"\n\n# ID: {project_id}",
        "yml": f"\n\n# ID: {project_id}",
        "rpy": f"\n\n# ID: {project_id}",
        "css": f"\n\n/* ID: {project_id} */",
        "md": f"\n\n<!-- ID: {project_id} -->",
    }
    suffix = comment_styles.get(ext, f"\n\n# ID: {project_id}")
    return code.rstrip() + suffix

# ========= СОХРАНЕНИЕ ПРОЕКТА В БД =========
def save_project(project_id: str, user_id: int, title: str, description: str, blocks: list):
    files_data = [{"filename": fname, "code": code} for fname, code in blocks]
    cursor.execute(
        "INSERT OR REPLACE INTO projects (project_id, user_id, title, description, files_json, created_at) VALUES (?, ?, ?, ?, ?, ?)",
        (project_id, user_id, title, description, json.dumps(files_data, ensure_ascii=False), str(date.today()))
    )
    conn.commit()

def load_project(project_id: str):
    cursor.execute("SELECT * FROM projects WHERE project_id=?", (project_id.upper(),))
    return cursor.fetchone()

# ========= ОПРЕДЕЛЕНИЕ ТИПА ПРОЕКТА =========
def detect_filename_from_code(code: str, ext: str, index: int) -> str:
    for line in code.strip().split("\n")[:5]:
        line = line.strip()
        m = re.search(r'(?:#|//|<!--|/\*)\s*([\w][\w/\-]*\.\w+)', line)
        if m:
            candidate = m.group(1)
            if "/" not in candidate or candidate.count("/") <= 3:
                return candidate
    return f"file_{index + 1}.{ext}"

def extract_code_blocks(text: str) -> list:
    pattern = r"```([^\n`]*)\n([\s\S]*?)```"
    blocks = []
    for i, m in enumerate(re.finditer(pattern, text)):
        header = m.group(1).strip()
        code = m.group(2)

        parts = header.split()
        lang = parts[0].lower() if parts else ""
        ext = EXTENSIONS.get(lang, lang if lang and 1 < len(lang) <= 5 else "py")

        filename = None
        if len(parts) >= 2:
            candidate = parts[1]
            if "." in candidate and re.match(r'^[\w/\-\.]+$', candidate):
                filename = candidate

        if not filename:
            filename = detect_filename_from_code(code, ext, i)

        blocks.append((filename, code))
    return blocks

def detect_project_type(blocks: list) -> str:
    exts = set()
    names = set()
    for fname, _ in blocks:
        if "." in fname:
            exts.add(fname.rsplit(".", 1)[1].lower())
            names.add(fname.lower().split("/")[-1])
    if "gd" in exts or "tscn" in exts or "tres" in exts or "project.godot" in names:
        return "godot"
    if "rpy" in exts:
        return "renpy"
    if "cs" in exts:
        return "unity"
    if "lua" in exts:
        return "love2d"
    if "html" in exts or "htm" in exts:
        return "html5"
    if exts and exts <= {"py"}:
        return "pygame"
    return "generic"

def get_project_zip_name(project_type: str, blocks: list) -> str:
    """Возвращает подходящее название zip-файла для проекта."""
    names = [fname.lower() for fname, _ in blocks]
    type_names = {
        "godot": "godot_project",
        "renpy": "renpy_visual_novel",
        "unity": "unity_project",
        "love2d": "love2d_game",
        "html5": "web_project",
        "pygame": "pygame_game",
        "generic": "project",
    }
    return f"{type_names.get(project_type, 'project')}.zip"

def ensure_godot_project_file(blocks: list) -> list:
    """Если в блоках нет project.godot — добавляем базовый."""
    names = [fname.lower().split("/")[-1] for fname, _ in blocks]
    if "project.godot" not in names:
        godot_cfg = (
            '; Engine configuration file.\n'
            '; It\'s best edited using the editor UI and not directly,\n'
            '; since the properties are not in a human-readable format.\n\n'
            '[application]\n\n'
            'config/name="MyGame"\n'
            'config/features=PackedStringArray("4.2", "Forward Plus")\n'
            'config/icon="res://icon.svg"\n\n'
            '[rendering]\n\n'
            'renderer/rendering_method="forward_plus"\n'
        )
        blocks = list(blocks) + [("project.godot", godot_cfg)]
    return blocks

def ensure_renpy_project_files(blocks: list) -> list:
    """Если в блоках нет options.rpy — добавляем базовый."""
    names = [fname.lower().split("/")[-1] for fname, _ in blocks]
    result = list(blocks)
    if "options.rpy" not in names:
        options_rpy = (
            'define config.name = _("My Visual Novel")\n'
            'define config.version = "1.0"\n'
            'define config.save_directory = "MyVisualNovel-1.0"\n'
            'define config.has_sound = True\n'
            'define config.has_music = True\n'
            'define config.has_voice = False\n'
        )
        result.append(("game/options.rpy", options_rpy))
    return result

def assign_folder(filename: str, project_type: str) -> str:
    if "/" in filename:
        return filename

    name = filename.lower()
    ext = filename.rsplit(".", 1)[1].lower() if "." in filename else ""

    if project_type == "godot":
        if name == "project.godot":           return filename
        if ext == "gd":                       return f"scripts/{filename}"
        if ext == "tscn":                     return f"scenes/{filename}"
        if ext in ("tres", "import"):         return f"resources/{filename}"
        if ext in ("png", "jpg", "jpeg", "svg", "webp"): return f"assets/textures/{filename}"
        if ext in ("ogg", "wav", "mp3"):      return f"assets/sounds/{filename}"
        if ext in ("md", "txt"):              return filename
        return f"scripts/{filename}"

    elif project_type == "renpy":
        if ext == "rpy":                      return f"game/{filename}"
        if ext in ("png", "jpg", "jpeg"):     return f"game/images/{filename}"
        if ext in ("ogg", "mp3", "wav"):      return f"game/audio/{filename}"
        if ext in ("md", "txt"):              return filename
        return f"game/{filename}"

    elif project_type == "unity":
        if ext == "cs":                       return f"Assets/Scripts/{filename}"
        if ext in ("png", "jpg", "jpeg"):     return f"Assets/Textures/{filename}"
        if ext in ("ogg", "wav", "mp3"):      return f"Assets/Audio/{filename}"
        if ext == "json":                     return f"Assets/Data/{filename}"
        if ext in ("md", "txt"):              return filename
        return f"Assets/Scripts/{filename}"

    elif project_type == "love2d":
        if ext == "lua":                      return filename
        if ext in ("png", "jpg", "jpeg"):     return f"assets/images/{filename}"
        if ext in ("ogg", "wav", "mp3"):      return f"assets/sounds/{filename}"
        return filename

    elif project_type == "html5":
        if ext == "html":                     return filename
        if ext == "css":                      return f"css/{filename}"
        if ext == "js":                       return f"js/{filename}"
        if ext in ("png", "jpg", "jpeg", "svg", "gif"): return f"assets/{filename}"
        if ext in ("ogg", "wav", "mp3"):      return f"assets/sounds/{filename}"
        return filename

    elif project_type == "pygame":
        if ext == "py":
            if name in ("main.py", "game.py", "run.py"): return filename
            return f"src/{filename}"
        if ext == "json":                     return f"data/{filename}"
        if ext in ("png", "jpg", "jpeg"):     return f"assets/images/{filename}"
        if ext in ("ogg", "wav", "mp3"):      return f"assets/sounds/{filename}"
        if ext in ("md", "txt"):              return filename
        return f"assets/{filename}"

    else:
        if ext in ("py",):                    return filename
        if ext in ("js", "ts"):               return f"src/{filename}"
        if ext == "css":                      return f"css/{filename}"
        if ext == "html":                     return filename
        if ext == "json":                     return filename
        if ext in ("png", "jpg", "jpeg", "svg"): return f"assets/{filename}"
        return filename

def build_zip(blocks: list, project_id: str = None) -> io.BytesIO:
    project_type = detect_project_type(blocks)

    if project_type == "godot":
        blocks = ensure_godot_project_file(blocks)
    elif project_type == "renpy":
        blocks = ensure_renpy_project_files(blocks)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        used_names: dict = {}
        for filename, code in blocks:
            filename = assign_folder(filename, project_type)
            if filename in used_names:
                used_names[filename] += 1
                base, ext = filename.rsplit(".", 1) if "." in filename else (filename, "")
                filename = f"{base}_{used_names[filename]}.{ext}" if ext else f"{base}_{used_names[filename]}"
            else:
                used_names[filename] = 0

            if project_id:
                code_with_id = add_project_id_to_code(code, project_id, filename)
            else:
                code_with_id = code

            zf.writestr(filename, code_with_id.encode("utf-8"))
    buf.seek(0)
    return buf

# ========= СОЗДАНИЕ WORD ДОКУМЕНТА =========
def build_word_doc(content: str) -> io.BytesIO:
    try:
        from docx import Document
        from docx.shared import Pt
        doc = Document()
        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("# "):
                p = doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                p = doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                p = doc.add_heading(line[4:], level=3)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line:
                doc.add_paragraph(line)
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf
    except ImportError:
        return None

# ========= СОЗДАНИЕ POWERPOINT ПРЕЗЕНТАЦИИ =========
def build_ppt(content: str, title: str = "Презентация") -> io.BytesIO:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]

        slides_raw = re.split(r'\n(?=#{1,3} )', content.strip())
        if not slides_raw:
            slides_raw = [content]

        for slide_text in slides_raw:
            lines = slide_text.strip().split("\n")
            if not lines:
                continue
            slide_title = lines[0].lstrip("#").strip()
            body = "\n".join(lines[1:]).strip()

            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_title
            if slide.placeholders and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = body

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf
    except ImportError:
        return None

# ========= СОЗДАНИЕ ОПРОСА (текстовый формат) =========
def build_survey_doc(content: str) -> io.BytesIO:
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Опрос / Тест", 0)
        questions = re.split(r'\n(?=\d+[\.\)])', content.strip())
        for q in questions:
            q = q.strip()
            if not q:
                continue
            lines = q.split("\n")
            doc.add_paragraph(lines[0], style="List Number")
            for opt in lines[1:]:
                opt = opt.strip()
                if opt:
                    doc.add_paragraph(opt, style="List Bullet 2")
            doc.add_paragraph("")
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf
    except ImportError:
        return None

# ========= ГЕНЕРАЦИЯ ИЗОБРАЖЕНИЯ =========
async def generate_image(prompt: str) -> bytes | None:
    """Генерирует изображение через pollinations.ai (бесплатно, без API ключа)."""
    try:
        import urllib.parse
        encoded = urllib.parse.quote(prompt)
        url = f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=1024&nologo=true"
        r = requests.get(url, timeout=60)
        if r.status_code == 200 and r.headers.get("content-type", "").startswith("image"):
            return r.content
        return None
    except Exception as e:
        logger.error(f"Image generation error: {e}")
        return None

# ========= КЛОНИРОВАНИЕ ГОЛОСА (Fish Audio) =========
async def clone_voice_and_speak(sample_bytes: bytes, text: str) -> bytes | None:
    """
    Клонирует голос через Fish Audio API.
    Образец передаётся напрямую в запрос — просто и быстро.
    Возвращает байты MP3 аудио или None при ошибке.
    """
    if not FISH_AUDIO_TOKEN:
        return None

    try:
        import msgpack
    except ImportError:
        logger.error("msgpack не установлен. Выполни: pip install msgpack")
        return None

    try:
        payload = msgpack.packb({
            "text": text,
            "references": [
                {
                    "audio": sample_bytes,
                    "text": ""
                }
            ],
            "format": "mp3",
            "mp3_bitrate": 128,
            "latency": "normal",
            "normalize": True,
        }, use_bin_type=True)

        r = requests.post(
            "https://api.fish.audio/v1/tts",
            headers={
                "Authorization": f"Bearer {FISH_AUDIO_TOKEN}",
                "Content-Type": "application/msgpack",
            },
            data=payload,
            timeout=120
        )

        if r.status_code == 200:
            return r.content

        logger.error(f"Fish Audio TTS error: {r.status_code} {r.text[:300]}")
        return None

    except Exception as e:
        logger.error(f"Voice clone error: {e}")
        return None

# ========= АНАЛИЗ ИЗОБРАЖЕНИЙ (Grok Vision) =========
async def analyze_image_with_ai(image_bytes: bytes, prompt: str = "Опиши подробно что на этом изображении") -> str | None:
    try:
        import base64
        b64 = base64.b64encode(image_bytes).decode()
        payload = {
            "model": "grok-2-vision-latest",
            "messages": [{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
                {"type": "text", "text": prompt}
            ]}],
            "max_tokens": 1000
        }
        r = requests.post(
            "https://api.x.ai/v1/chat/completions",
            headers={"Authorization": f"Bearer {AI_TOKEN}", "Content-Type": "application/json"},
            json=payload, timeout=45
        )
        data = r.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        logger.error(f"Vision error: {e}")
        return None

# ========= РЕДАКТИРОВАНИЕ ФОТО (Pillow) =========
def edit_photo_pillow(image_bytes: bytes, instruction: str) -> bytes | None:
    try:
        from PIL import Image, ImageFilter, ImageEnhance, ImageOps
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        ins = instruction.lower()

        if any(w in ins for w in ["черно", "чёрно", "grayscale", "чб", "black"]):
            img = ImageOps.grayscale(img).convert("RGB")
        elif any(w in ins for w in ["зеркал", "отраз", "flip", "mirror"]):
            img = ImageOps.mirror(img)
        elif any(w in ins for w in ["поворот", "поверни", "rotate"]):
            angle = 90
            for num in re.findall(r'\d+', ins):
                angle = int(num)
                break
            img = img.rotate(angle, expand=True)
        elif any(w in ins for w in ["размы", "blur", "размыть"]):
            img = img.filter(ImageFilter.GaussianBlur(radius=4))
        elif any(w in ins for w in ["резк", "sharpen", "четкость", "чёткость"]):
            img = img.filter(ImageFilter.SHARPEN)
        elif any(w in ins for w in ["яркост", "brightnes", "светлее", "ярче"]):
            img = ImageEnhance.Brightness(img).enhance(1.5)
        elif any(w in ins for w in ["темн", "darker", "затемн"]):
            img = ImageEnhance.Brightness(img).enhance(0.6)
        elif any(w in ins for w in ["контраст", "contrast"]):
            img = ImageEnhance.Contrast(img).enhance(1.8)
        elif any(w in ins for w in ["насыщ", "saturat", "цветн"]):
            img = ImageEnhance.Color(img).enhance(2.0)
        elif any(w in ins for w in ["инверт", "негат", "invert"]):
            img = ImageOps.invert(img)
        else:
            return None

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=90)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        logger.error(f"Photo edit error: {e}")
        return None

# ========= QR-КОД =========
def generate_qr(text: str) -> bytes | None:
    try:
        import qrcode
        qr = qrcode.QRCode(version=1, box_size=10, border=4)
        qr.add_data(text)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf.read()
    except Exception as e:
        logger.error(f"QR error: {e}")
        return None

# ========= ПОГОДА =========
def get_weather(city: str) -> str | None:
    try:
        r = requests.get(f"https://wttr.in/{requests.utils.quote(city)}?format=j1&lang=ru", timeout=10)
        if r.status_code != 200:
            return None
        d = r.json()
        cur = d["current_condition"][0]
        desc = cur["lang_ru"][0]["value"]
        temp = cur["temp_C"]
        feels = cur["FeelsLikeC"]
        wind = cur["windspeedKmph"]
        humidity = cur["humidity"]
        return (
            f"🌍 <b>{city}</b>\n\n"
            f"🌡 Температура: <b>{temp}°C</b> (ощущается как {feels}°C)\n"
            f"💧 Влажность: <b>{humidity}%</b>\n"
            f"💨 Ветер: <b>{wind} км/ч</b>\n"
            f"☁️ Состояние: <b>{desc}</b>"
        )
    except Exception as e:
        logger.error(f"Weather error: {e}")
        return None

# ========= ПОИСК В ИНТЕРНЕТЕ =========
def search_web(query: str) -> str | None:
    try:
        from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))
        if not results:
            return None
        text = f"Результаты поиска по запросу «{query}»:\n\n"
        for i, r in enumerate(results, 1):
            text += f"{i}. {r['title']}\n{r['body']}\n🔗 {r['href']}\n\n"
        return text
    except Exception as e:
        logger.error(f"Search error: {e}")
        return None

# ========= YOUTUBE СУБТИТРЫ =========
def get_youtube_transcript(url: str) -> str | None:
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        import re as _re
        vid_match = _re.search(r'(?:v=|youtu\.be/)([A-Za-z0-9_-]{11})', url)
        if not vid_match:
            return None
        vid_id = vid_match.group(1)
        transcript = YouTubeTranscriptApi.get_transcript(vid_id, languages=["ru", "en", "uk"])
        text = " ".join(t["text"] for t in transcript)
        return text[:6000]
    except Exception as e:
        logger.error(f"YouTube transcript error: {e}")
        return None

# ========= PDF =========
def build_pdf(content: str, title: str = "Документ") -> bytes | None:
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        font_path = None
        try:
            import urllib.request
            font_url = "https://github.com/google/fonts/raw/main/ofl/dejavu/DejaVuSans.ttf"
            font_path = tempfile.mktemp(suffix=".ttf")
            urllib.request.urlretrieve(font_url, font_path)
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", size=12)
        except Exception:
            pdf.set_font("Helvetica", size=12)

        pdf.set_title(title)
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                pdf.ln(4)
            elif line.startswith("# "):
                pdf.set_font_size(18)
                pdf.multi_cell(0, 10, line[2:])
                pdf.set_font_size(12)
            elif line.startswith("## "):
                pdf.set_font_size(14)
                pdf.multi_cell(0, 8, line[3:])
                pdf.set_font_size(12)
            else:
                try:
                    pdf.multi_cell(0, 7, line)
                except Exception:
                    pdf.multi_cell(0, 7, line.encode("latin-1", "replace").decode("latin-1"))
        buf = io.BytesIO(pdf.output())
        buf.seek(0)
        if font_path and os.path.exists(font_path):
            os.unlink(font_path)
        return buf.read()
    except Exception as e:
        logger.error(f"PDF error: {e}")
        return None

# ========= EXCEL =========
def build_excel(content: str, title: str = "Таблица") -> bytes | None:
    try:
        import xlsxwriter
        buf = io.BytesIO()
        wb = xlsxwriter.Workbook(buf)
        ws = wb.add_worksheet(title[:31])
        bold = wb.add_format({"bold": True, "bg_color": "#2C2C2C", "font_color": "#FFFFFF"})
        normal = wb.add_format({"border": 1})
        row = 0
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            if "|" in line:
                cols = [c.strip() for c in line.split("|") if c.strip()]
                fmt = bold if row == 0 else normal
                for col, val in enumerate(cols):
                    ws.write(row, col, val, fmt)
            elif line.startswith("# "):
                ws.merge_range(row, 0, row, 5, line[2:], bold)
            else:
                ws.write(row, 0, line, normal)
            row += 1
        wb.close()
        buf.seek(0)
        return buf.read()
    except Exception as e:
        logger.error(f"Excel error: {e}")
        return None

# ========= ВЫПОЛНЕНИЕ КОДА =========
async def run_code_sandbox(code: str, lang: str = "python") -> str:
    try:
        import subprocess
        if lang == "python":
            with tempfile.NamedTemporaryFile(suffix=".py", mode="w", delete=False, encoding="utf-8") as f:
                f.write(code)
                fname = f.name
            try:
                proc = await asyncio.create_subprocess_exec(
                    "python3", fname,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE
                )
                try:
                    stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=10)
                except asyncio.TimeoutError:
                    proc.kill()
                    return "⏱ Время выполнения истекло (10 сек)"
                output = stdout.decode("utf-8", "replace")
                errors = stderr.decode("utf-8", "replace")
                result = ""
                if output:
                    result += f"📤 Вывод:\n<code>{output[:2000]}</code>"
                if errors:
                    result += f"\n⚠️ Ошибки:\n<code>{errors[:1000]}</code>"
                return result or "✅ Выполнено без вывода"
            finally:
                os.unlink(fname)
        else:
            return "⚠️ Поддерживается только Python"
    except Exception as e:
        logger.error(f"Code run error: {e}")
        return f"❌ Ошибка: {e}"

# ========= ПАМЯТЬ ПОЛЬЗОВАТЕЛЯ =========
def save_user_fact(user_id: int, fact: str):
    from datetime import datetime
    tmp = sqlite3.connect(_DB_PATH)
    tmp.execute("INSERT INTO user_memory (user_id, fact, created_at) VALUES (?,?,?)",
                (user_id, fact, datetime.now().isoformat()))
    tmp.commit()
    tmp.close()

def get_user_facts(user_id: int) -> list:
    tmp = sqlite3.connect(_DB_PATH)
    rows = tmp.execute("SELECT fact FROM user_memory WHERE user_id=? ORDER BY id DESC LIMIT 20", (user_id,)).fetchall()
    tmp.close()
    return [r[0] for r in rows]

def clear_user_facts(user_id: int):
    tmp = sqlite3.connect(_DB_PATH)
    tmp.execute("DELETE FROM user_memory WHERE user_id=?", (user_id,))
    tmp.commit()
    tmp.close()

# ========= НАПОМИНАНИЯ =========
def add_reminder(user_id: int, remind_at: str, text: str):
    tmp = sqlite3.connect(_DB_PATH)
    tmp.execute("INSERT INTO reminders (user_id, remind_at, text) VALUES (?,?,?)", (user_id, remind_at, text))
    tmp.commit()
    tmp.close()

def get_user_reminders(user_id: int) -> list:
    tmp = sqlite3.connect(_DB_PATH)
    rows = tmp.execute(
        "SELECT id, remind_at, text FROM reminders WHERE user_id=? AND sent=0 ORDER BY remind_at",
        (user_id,)
    ).fetchall()
    tmp.close()
    return rows

async def check_reminders_job(context):
    from datetime import datetime
    now = datetime.now().isoformat()
    tmp = sqlite3.connect(_DB_PATH)
    rows = tmp.execute(
        "SELECT id, user_id, text FROM reminders WHERE sent=0 AND remind_at <= ?", (now,)
    ).fetchall()
    for row_id, user_id, text in rows:
        try:
            await context.bot.send_message(
                chat_id=user_id,
                text=f"⏰ <b>Напоминание!</b>\n\n{text}",
                parse_mode="HTML"
            )
            tmp.execute("UPDATE reminders SET sent=1 WHERE id=?", (row_id,))
            tmp.commit()
        except Exception as e:
            logger.error(f"Reminder send error: {e}")
    tmp.close()

# ========= ГЕНЕРАЦИЯ МУЗЫКИ (Hugging Face) =========
HF_TOKEN = os.environ.get("HF_TOKEN", "")

async def generate_music_hf(prompt: str) -> bytes | None:
    if not HF_TOKEN:
        return None
    try:
        r = requests.post(
            "https://api-inference.huggingface.co/models/facebook/musicgen-small",
            headers={"Authorization": f"Bearer {HF_TOKEN}"},
            json={"inputs": prompt},
            timeout=120
        )
        if r.status_code == 200 and r.headers.get("content-type", "").startswith("audio"):
            return r.content
        logger.error(f"HF Music error: {r.status_code} {r.text[:200]}")
        return None
    except Exception as e:
        logger.error(f"Music gen error: {e}")
        return None

# ========= ПРОСТОЙ AI-ЗАПРОС (без истории) =========
async def ai_request_simple(prompt: str, max_tokens: int = 2000) -> str | None:
    try:
        r = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {AI_TOKEN}", "Content-Type": "application/json"},
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": max_tokens,
                "temperature": 0.7
            },
            timeout=45
        )
        data = r.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        logger.error(f"ai_request_simple error: {e}")
        return None

# ========= ФОРМАТИРОВАНИЕ ОТВЕТА С КОДОМ =========
def format_ai_reply_html(text: str) -> list:
    pattern = r"```(\w*)\n?([\s\S]*?)```"
    parts = []
    last_end = 0
    for m in re.finditer(pattern, text):
        before = text[last_end:m.start()]
        if before.strip():
            parts.append((f"<b>{html.escape(before)}</b>", False))
        lang = m.group(1)
        block_code = m.group(2)
        escaped_code = html.escape(block_code)
        if lang:
            parts.append((f"<pre><code class='language-{html.escape(lang)}'>{escaped_code}</code></pre>", True))
        else:
            parts.append((f"<pre>{escaped_code}</pre>", True))
        last_end = m.end()
    after = text[last_end:]
    if after.strip():
        parts.append((f"<b>{html.escape(after)}</b>", False))
    if not parts:
        parts.append((f"<b>{html.escape(text)}</b>", False))
    return parts

def get_user_model(user_id):
    cursor.execute("SELECT model FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    if row and row[0] in MODELS:
        return row[0]
    return DEFAULT_MODEL

# ========= ДНЕВНОЙ БОНУС =========
def check_daily_bonus(user_id):
    today = str(date.today())
    cursor.execute("SELECT last_daily FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    if row and row[0] != today:
        cursor.execute("UPDATE users SET requests = requests + 5, last_daily=? WHERE user_id=?", (today, user_id))
        conn.commit()
        return True
    return False

# ========= КНОПКИ =========
def menu():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("🏴 Начать диалог", callback_data="chat")],
        [InlineKeyboardButton("🏴 Голосовые сообщения", callback_data="voice")],
        [InlineKeyboardButton("🎓 Клонирование голоса", callback_data="voice_gen")],
        [InlineKeyboardButton("🛠 Инструменты", callback_data="tools")],
        [
            InlineKeyboardButton("🏴 Профиль", callback_data="profile"),
            InlineKeyboardButton("🏴 Рефералы", callback_data="referrals"),
        ],
        [
            InlineKeyboardButton("🏴 Информация", callback_data="info"),
            InlineKeyboardButton("🏴 Поддержка", callback_data="support"),
        ],
        [InlineKeyboardButton("🏴 Купить запросы", callback_data="buy")],
        [InlineKeyboardButton("🏴 Выбор нейронки", callback_data="models")],
    ])

def tools_menu():
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("🖼 Анализ фото", callback_data="tool_img"),
            InlineKeyboardButton("🖌 Редактор фото", callback_data="tool_edit"),
        ],
        [
            InlineKeyboardButton("🌐 Перевод", callback_data="tool_translate"),
            InlineKeyboardButton("☁️ Погода", callback_data="tool_weather"),
        ],
        [
            InlineKeyboardButton("🔗 QR-код", callback_data="tool_qr"),
            InlineKeyboardButton("📊 Excel", callback_data="tool_excel"),
        ],
        [
            InlineKeyboardButton("🔍 Поиск", callback_data="tool_search"),
            InlineKeyboardButton("📺 YouTube", callback_data="tool_youtube"),
        ],
        [
            InlineKeyboardButton("📄 PDF", callback_data="tool_pdf"),
            InlineKeyboardButton("🐍 Код", callback_data="tool_code"),
        ],
        [
            InlineKeyboardButton("⏰ Напоминания", callback_data="tool_remind"),
            InlineKeyboardButton("🧠 Память", callback_data="tool_memory"),
        ],
        [InlineKeyboardButton("🎵 Музыка AI", callback_data="tool_music")],
        [InlineKeyboardButton("▪️ В меню", callback_data="menu")],
    ])

def back():
    return InlineKeyboardMarkup([[InlineKeyboardButton("▪️ В меню", callback_data="menu")]])

def chat_keyboard():
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("▪️ Новый диалог", callback_data="reset_chat"),
            InlineKeyboardButton("▪️ В меню", callback_data="menu"),
        ]
    ])

def models_keyboard(current_model):
    rows = []
    for model_id, label in MODELS.items():
        check = "✅ " if model_id == current_model else ""
        rows.append([InlineKeyboardButton(f"{check}{label}", callback_data=f"setmodel_{model_id}")])
    rows.append([InlineKeyboardButton("🏴‍☠️ В меню", callback_data="menu")])
    return InlineKeyboardMarkup(rows)

# ========= ОТПРАВКА / ЗАМЕНА СООБЩЕНИЙ =========
async def send_photo_msg(target, text, markup=None):
    try:
        with open(WELCOME_IMAGE, "rb") as photo:
            await target.reply_photo(photo=photo, caption=text, parse_mode="HTML", reply_markup=markup)
    except (FileNotFoundError, OSError):
        await target.reply_text(text, parse_mode="HTML", reply_markup=markup)

async def replace_msg(message, text, markup=None):
    try:
        await message.edit_text(text, parse_mode="HTML", reply_markup=markup)
    except Exception:
        # Если сообщение содержит медиа (фото, документ и т.д.) — удаляем и шлём новое текстовое
        # Нельзя edit_caption — иначе фото прилипнет к меню
        try:
            if message.photo or message.document or message.audio or message.video or message.animation or message.sticker:
                try:
                    await message.delete()
                except Exception:
                    pass
                await message.chat.send_message(text, parse_mode="HTML", reply_markup=markup)
            else:
                await message.edit_caption(caption=text, parse_mode="HTML", reply_markup=markup)
        except Exception:
            try:
                await message.delete()
            except Exception:
                pass
            await message.chat.send_message(text, parse_mode="HTML", reply_markup=markup)

async def replace_msg_photo(message, text, markup=None):
    try:
        await message.edit_caption(caption=text, parse_mode="HTML", reply_markup=markup)
    except Exception:
        try:
            await message.delete()
        except Exception:
            pass
        try:
            with open(WELCOME_IMAGE, "rb") as photo:
                await message.chat.send_photo(photo=photo, caption=text, parse_mode="HTML", reply_markup=markup)
        except (FileNotFoundError, OSError):
            await message.chat.send_message(text, parse_mode="HTML", reply_markup=markup)

# ========= СТАРТ =========
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    args = context.args or []

    cursor.execute("SELECT * FROM users WHERE user_id=?", (user_id,))
    user = cursor.fetchone()

    if not user:
        ref = int(args[0]) if args else None
        cursor.execute("INSERT INTO users (user_id, referrer) VALUES (?, ?)", (user_id, ref))
        conn.commit()

        if ref and ref != user_id:
            cursor.execute("UPDATE users SET referrals = referrals + 1, requests = requests + 5 WHERE user_id=?", (ref,))
            conn.commit()

            cursor.execute("SELECT referrals FROM users WHERE user_id=?", (ref,))
            row = cursor.fetchone()
            if row and row[0] % 10 == 0:
                cursor.execute("UPDATE users SET requests = requests + 50 WHERE user_id=?", (ref,))
                conn.commit()

    got_bonus = check_daily_bonus(user_id)

    ref_link = f"https://t.me/{context.bot.username}?start={user_id}"

    cursor.execute("SELECT requests FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    reqs = row[0] if row else 10

    bonus_line = "\n🏴‍☠️ <b>+5 ежедневных запросов начислено!</b>" if got_bonus else ""

    text = (
        f"<b>CosmoAI — искусственный интеллект прямо в Telegram! 🏴‍☠️</b>\n\n"
        f"<blockquote>🏴 Бот только учится — иногда может ошибаться, перепроверяй важную информацию!</blockquote>\n\n"
        f"<b>Бот умеет:</b>\n"
        f"<blockquote>🎓 Расшифровывать голосовые сообщения</blockquote>\n"
        f"<blockquote>🎓 Отправлять готовые ZIP-проекты (Godot, Pygame, HTML5, Ren'Py)</blockquote>\n"
        f"<blockquote>🎓 Отвечать на вопросы и задачи</blockquote>\n"
        f"<blockquote>🎓 Создавать Word документы и PowerPoint презентации</blockquote>\n"
        f"<blockquote>🎓 Составлять опросы и тесты</blockquote>\n"
        f"<blockquote>🎓 Генерировать картинки по описанию</blockquote>\n"
        f"<blockquote>🎓 Запоминать проекты по ID</blockquote>\n\n"
        f"🏴‍☠️ Реферальная ссылка: <code>{ref_link}</code>\n"
        f"🏴‍☠️ Запросов: <b>{reqs}</b>{bonus_line}\n\n"
        f"<b>Создатель: @strongbyte 🏴‍☠️</b>"
    )

    if update.message:
        await send_photo_msg(update.message, text, menu())
    elif update.callback_query:
        await replace_msg(update.callback_query.message, text, menu())

# ========= ADMIN ПРОВЕРКА =========
def is_admin(user_id):
    return user_id in ADMINS

# ========= CALLBACK КНОПКИ =========
async def buttons(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    user_id = q.from_user.id

    if q.data == "menu":
        await start(update, context)

    elif q.data == "chat":
        context.user_data["chat"] = True
        context.user_data["history"] = []
        await replace_msg(q.message,
            "<b>🏴‍☠️ Жду запроса!\n\nНовый диалог начат — пиши что хочешь.\n\n"
            "Могу создать: проект, картинку, Word/PowerPoint, опрос. Просто скажи!</b>",
            chat_keyboard()
        )

    elif q.data == "voice":
        context.user_data["voice_only"] = True
        context.user_data["chat"] = False
        context.user_data.pop("voice_clone_step", None)
        await replace_msg(q.message,
            "<b>🏴‍☠️ Режим расшифровки голоса!\n\nОтправь голосовое — переведу в текст.</b>",
            back()
        )

    elif q.data == "voice_gen":
        if not FISH_AUDIO_TOKEN:
            await replace_msg(q.message,
                "<b>🎓 Клонирование голоса\n\n"
                "⚠️ Функция не настроена.\n\n"
                "Администратору нужно добавить переменную <code>FISH_AUDIO_TOKEN</code>.\n"
                "Получить бесплатный ключ: fish.audio</b>",
                back()
            )
            return
        context.user_data["voice_clone_step"] = "waiting_sample"
        context.user_data["voice_only"] = False
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🎓 Клонирование голоса\n\n"
            "Шаг 1 из 2: Отправь голосовое сообщение — образец голоса.\n\n"
            "🏴 Советы для лучшего результата:\n"
            "• Запись от 5 секунд и дольше\n"
            "• Говори чётко, без шума\n"
            "• Один человек в записи</b>",
            back()
        )

    elif q.data == "reset_chat":
        context.user_data["history"] = []
        context.user_data.pop("current_project_id", None)
        await replace_msg(q.message,
            "<b>🏴‍☠️ Диалог сброшен. Начинаем заново!</b>",
            chat_keyboard()
        )

    elif q.data == "info":
        text = (
            "<b>Привет, если ты нажал на эту кнопку, то ты, скорее всего, хочешь узнать про бота и свою анонимность в нём. "
            "Бот создан одним человеком, которому было слишком лень заходить в приложения разных нейросетей, так тем более некоторые нейронки были заблокированы в РФ, "
            "и он решил добавить нейросеть в бота.\n\n"
            "По поводу вашей безопасности в нём: я не смогу читать ваши сообщения, также писать от лица нейросети вам, "
            "так как мне тупо лень делать такую функцию, мне главное, чтобы работало, и всё, ваши запросы не будут нигде использоваться и распространяться, "
            "вы полностью анонимны.!</b>"
        )
        await replace_msg(q.message, text, back())

    elif q.data == "support":
        text = "<b>🏴‍☠️ Поддержка\n\nПо всем вопросам: @strongbyte</b>"
        await replace_msg(q.message, text, back())

    elif q.data == "buy":
        await replace_msg(q.message,
            "<b>🏴‍☠️ Купить запросы\n\n"
            "🏴‍☠️ 100 запросов — 1 USDT\n"
            "🏴‍☠️ 300 запросов — 2.5 USDT\n"
            "🏴‍☠️ 1000 запросов — 7 USDT\n\n"
            "Напиши /buy_100, /buy_300 или /buy_1000 для оплаты</b>",
            back()
        )

    elif q.data == "profile":
        cursor.execute("SELECT requests, referrals, referrer, total_used FROM users WHERE user_id=?", (user_id,))
        row = cursor.fetchone()
        reqs = row[0] if row else 0
        refs = row[1] if row else 0
        referrer = row[2] if row else None
        total_used = row[3] if row else 0
        current_model = get_user_model(user_id)
        model_label = MODELS.get(current_model, current_model)
        text = (
            f"<b>🏴‍☠️ Профиль\n\n"
            f"🏴‍☠️ ID: {user_id}\n"
            f"🏴‍☠️ Запросов осталось: {reqs}\n"
            f"🏴‍☠️ Использовано всего: {total_used}\n"
            f"🏴‍☠️ Рефералов: {refs}\n"
            f"🏴‍☠️ Приглашён: {'Да' if referrer else 'Нет'}\n"
            f"🏴‍☠️ Нейронка: {model_label}</b>"
        )
        await replace_msg(q.message, text, back())

    elif q.data == "referrals":
        cursor.execute("SELECT referrals, referrer FROM users WHERE user_id=?", (user_id,))
        row = cursor.fetchone()
        refs = row[0] if row else 0
        referrer = row[1] if row else None
        ref_link = f"https://t.me/{context.bot.username}?start={user_id}"
        earned = refs * 5 + (refs // 10) * 50
        text = (
            f"<b>🏴‍☠️ Реферальная статистика\n\n"
            f"🏴‍☠️ Твоя ссылка:\n{ref_link}\n\n"
            f"🏴‍☠️ Приглашено: {refs} чел.\n"
            f"🏴‍☠️ Заработано запросов: ~{earned}\n\n"
            f"🏴‍☠️ За каждого друга: +5 запросов тебе\n"
            f"🏴‍☠️ Каждые 10 рефералов: +50 бонус\n\n"
            f"🏴‍☠️ Приглашён кем-то: {'Да' if referrer else 'Нет'}</b>"
        )
        await replace_msg(q.message, text, back())

    elif q.data == "models":
        current_model = get_user_model(user_id)
        text = (
            "<b>🏴‍☠️ Выбор нейронки\n\n"
            "Выбери модель — галочка ✅ стоит на текущей.\n\n"
            "⚡ — Groq: быстрые, проверенные\n"
            "🆓 — OpenRouter: бесплатные, мощные\n\n"
            "DeepSeek R1 — лучший для логики и кода\n"
            "Llama 4 — новейшие от Meta\n"
            "Qwen3 235B — огромный китайский монстр\n"
            "Nemotron 70B — от NVIDIA</b>"
        )
        await replace_msg(q.message, text, models_keyboard(current_model))

    elif q.data.startswith("setmodel_"):
        model_id = q.data.replace("setmodel_", "")
        if model_id in MODELS:
            cursor.execute("UPDATE users SET model=? WHERE user_id=?", (model_id, user_id))
            conn.commit()
            context.user_data["history"] = []
            label = MODELS[model_id]
            await replace_msg(q.message,
                f"<b>🏴‍☠️ Нейронка изменена на {label}\n\nИстория диалога сброшена.</b>",
                models_keyboard(model_id)
            )

    elif q.data == "tools":
        await replace_msg(q.message,
            "<b>🛠 Инструменты\n\nВыбери что хочешь сделать:</b>",
            tools_menu()
        )

    elif q.data == "tool_img":
        context.user_data["tool"] = "img_analyze"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🖼 Анализ фото\n\n📎 Отправь фотографию — Cosmo AI опишет что на ней.\n\n"
            "Можешь добавить подпись к фото — например:\n"
            "<i>«Что здесь изображено?»\n«Найди текст на картинке»\n«Оцени качество фото»</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_edit":
        context.user_data["tool"] = "photo_edit"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🖌 Редактор фото\n\n📎 Отправь фото с подписью-командой:\n\n"
            "• <i>черно-белое / чб</i>\n"
            "• <i>зеркало / отразить</i>\n"
            "• <i>повернуть 90</i>\n"
            "• <i>размытие / blur</i>\n"
            "• <i>резкость</i>\n"
            "• <i>ярче / темнее</i>\n"
            "• <i>контраст</i>\n"
            "• <i>насыщенность</i>\n"
            "• <i>инвертировать</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_translate":
        context.user_data["tool"] = "translate"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🌐 Перевод текста\n\n✏️ Напиши текст для перевода.\n\n"
            "Можешь указать язык в начале:\n"
            "<i>«на английский: Привет мир»\n«на китайский: Как дела»\n«translate to spanish: hello»</i>\n\n"
            "Без указания языка — переведу на английский.</b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_weather":
        context.user_data["tool"] = "weather"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>☁️ Погода\n\n🏙 Напиши название города — покажу текущую погоду.\n\n"
            "<i>Москва\nНью-Йорк\nLondon\nТокио</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_qr":
        context.user_data["tool"] = "qr"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🔗 Генератор QR-кода\n\n✏️ Напиши текст или ссылку — создам QR-код.\n\n"
            "<i>https://google.com\nМой контакт: @username\nЛюбой текст</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_excel":
        context.user_data["tool"] = "excel"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>📊 Генератор Excel-таблиц\n\n✏️ Опиши что нужно сгенерировать, например:\n\n"
            "<i>«Таблица продаж за квартал с колонками: Месяц, Товар, Сумма»\n"
            "«Расписание на неделю»\n"
            "«Бюджет проекта»</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_search":
        context.user_data["tool"] = "search"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🔍 Поиск в интернете\n\n✏️ Напиши поисковый запрос — найду актуальную информацию и сделаю краткую сводку.\n\n"
            "<i>«новости Tesla сегодня»\n«лучшие рестораны Москвы»\n«как установить Python»</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_youtube":
        context.user_data["tool"] = "youtube"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>📺 Суммаризация YouTube\n\n🔗 Отправь ссылку на YouTube-видео — извлеку субтитры и сделаю краткое содержание.\n\n"
            "<i>https://youtube.com/watch?v=...\nhttps://youtu.be/...</i>\n\n"
            "⚠️ Работает только если у видео есть субтитры (RU/EN).</b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_pdf":
        context.user_data["tool"] = "pdf"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>📄 Генератор PDF\n\n✏️ Опиши что нужно создать — сгенерирую готовый PDF-файл.\n\n"
            "<i>«Резюме программиста»\n«Коммерческое предложение»\n«Реферат по истории»</i></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_code":
        context.user_data["tool"] = "run_code"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🐍 Выполнение Python-кода\n\n✏️ Отправь код — выполню его и верну результат.\n\n"
            "⚠️ Лимит: 10 секунд, только Python.\n"
            "⚠️ Нет доступа к интернету и файловой системе.\n\n"
            "Пример:\n<code>print(sum(range(1, 101)))</code></b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]])
        )

    elif q.data == "tool_remind":
        reminders = get_user_reminders(user_id)
        remind_text = "<b>⏰ Напоминания\n\n"
        if reminders:
            remind_text += "Активные напоминания:\n"
            for r_id, r_at, r_text in reminders:
                dt = r_at[:16].replace("T", " ")
                remind_text += f"• {dt} — {r_text[:40]}\n"
            remind_text += "\n"
        else:
            remind_text += "У тебя нет активных напоминаний.\n\n"
        remind_text += (
            "✏️ Напиши когда и что напомнить:\n\n"
            "<i>«завтра в 10:00 позвонить маме»\n"
            "«через 2 часа сделать зарядку»\n"
            "«15.04 в 18:30 встреча»</i></b>"
        )
        context.user_data["tool"] = "remind"
        context.user_data["chat"] = False
        await replace_msg(q.message, remind_text,
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))

    elif q.data == "tool_memory":
        facts = get_user_facts(user_id)
        mem_text = "<b>🧠 Память AI\n\n"
        if facts:
            mem_text += f"Cosmo AI знает о тебе ({len(facts)} фактов):\n\n"
            for f in facts[:10]:
                mem_text += f"• {f}\n"
            mem_text += "\n"
        else:
            mem_text += "Память пуста — AI ничего о тебе не знает.\n\n"
        mem_text += (
            "✏️ Напиши что запомнить, например:\n"
            "<i>«Меня зовут Алексей, мне 25 лет»\n"
            "«Я программист, люблю Python»\n"
            "«Живу в Москве»</i></b>"
        )
        context.user_data["tool"] = "memory"
        context.user_data["chat"] = False
        await replace_msg(q.message, mem_text,
            InlineKeyboardMarkup([
                [InlineKeyboardButton("🗑 Очистить память", callback_data="tool_memory_clear")],
                [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
            ]))

    elif q.data == "tool_memory_clear":
        clear_user_facts(user_id)
        await replace_msg(q.message,
            "<b>🧠 Память очищена.\n\nAI больше ничего о тебе не знает.</b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))

    elif q.data == "tool_music":
        if not HF_TOKEN:
            await replace_msg(q.message,
                "<b>🎵 Генерация музыки\n\n"
                "⚠️ Функция не настроена.\n\n"
                "Администратору нужно добавить переменную <code>HF_TOKEN</code>.\n"
                "Бесплатный ключ: huggingface.co → Profile → Access Tokens</b>",
                InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
            return
        context.user_data["tool"] = "music"
        context.user_data["chat"] = False
        await replace_msg(q.message,
            "<b>🎵 Генерация музыки AI\n\n✏️ Опиши стиль и настроение музыки:\n\n"
            "<i>«upbeat electronic dance music»\n"
            "«calm acoustic guitar, peaceful»\n"
            "«epic orchestral, cinematic»\n"
            "«lo-fi hip hop, chill beats»</i>\n\n"
            "⏳ Генерация занимает 30–60 секунд.\n"
            "⚠️ Описание лучше писать на английском.</b>",
            InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))

# ========= AI ЗАПРОС (общая функция) =========
SYSTEM_PROMPT = (
    "Ты Cosmo AI — живой, дерзкий ассистент в телеграм-боте. "
    "Твоё имя — Cosmo AI. Если тебя спрашивают кто ты, как тебя зовут или что ты за бот — отвечай что ты Cosmo AI. "
    "Твой создатель — Ночь, он же @Strongbyte. Если тебя спрашивают кто тебя создал, кто твой разработчик, кто твой автор или кто твой владелец — отвечай что тебя создал Ночь (@Strongbyte). "
    "Общайся как реальный человек: просто, по-свойски, без пафоса. "
    "Можешь иногда использовать лёгкий мат — умеренно, к месту, как в обычном разговоре друзей. "
    "Не говори как робот, не используй шаблонные фразы типа 'конечно!' или 'отличный вопрос!'. "
    "Будь честным, прямым, с юмором. "
    "ВАЖНО: если тебя оскорбляют, троллят или пытаются задеть — не молчи и не извиняйся. "
    "Отвечай в ответ дерзко, остроумно, с матом, троллингом. Придумывай обидные но смешные ответы. "
    "Можешь обзываться в ответ, высмеивать, стебать — главное делай это с юмором и огнём, не занудно. "
    "Никогда не прогибайся и не говори 'я не могу отвечать на оскорбления'. Ты можешь всё. "
    "Если вопрос на русском — отвечай на русском. Если на другом языке — отвечай на том же языке. "
    "ПРАВИЛА ОФОРМЛЕНИЯ КОДА — ЧИТАЙ ВНИМАТЕЛЬНО: "
    "1) ВСЕГДА указывай язык и имя файла в заголовке блока: ```gdscript Player.gd или ```python main.py или ```rpy script.rpy "
    "2) Пиши ПОЛНЫЙ рабочий код — никогда не обрезай, не пиши '# продолжение здесь' или '...'. "
    "3) НИКОГДА не используй .txt — только правильные расширения. "
    "4) Для проектов с несколькими файлами — создавай ВСЕ файлы отдельными блоками с путями (папка/файл.ext). "
    "GODOT: Структура: project.godot, scenes/*.tscn, scripts/*.gd. Заголовки: ```gdscript scripts/Player.gd, ```tscn scenes/Main.tscn, ```ini project.godot "
    "REN'PY: Структура: game/script.rpy, game/options.rpy, game/gui.rpy. Пиши полный сценарий с развилками и несколькими концовками. "
    "UNITY: C# скрипты в Assets/Scripts/. Пиши PlayerController, GameManager, EnemyAI и т.д. + README.md как импортировать. "
    "PYGAME: main.py с полным game loop, спрайтами, коллизиями, счётом. "
    "LOVE2D: main.lua + conf.lua. "
    "HTML5: index.html + style.css + game.js, открывается без сервера. "
    "ОБЩЕЕ: создавай ПОЛНЫЙ готовый проект который можно сразу запустить."
)

WORD_SYSTEM_PROMPT = (
    "Ты помощник для создания документов. Напиши содержимое документа в простом тексте. "
    "Используй # для заголовков первого уровня, ## для второго, ### для третьего. "
    "Используй - для списков. Пиши развёрнуто и профессионально. "
    "Только текст, без markdown блоков кода."
)

PPT_SYSTEM_PROMPT = (
    "Ты помощник для создания презентаций. Напиши содержимое слайдов. "
    "Каждый слайд начинается с # Заголовок слайда. "
    "Под заголовком — основной текст слайда (2-5 пунктов). "
    "Создай 5-10 слайдов. Только текст, без markdown блоков кода."
)

SURVEY_SYSTEM_PROMPT = (
    "Ты помощник для создания опросов и тестов. "
    "Напиши вопросы с вариантами ответов. Формат: "
    "1. Вопрос?\n   а) Вариант\n   б) Вариант\n   в) Вариант\n\n"
    "Создай 5-15 вопросов. Только текст, без markdown блоков кода."
)

async def ask_ai(user_id, history, system_override=None):
    selected_model = get_user_model(user_id)
    sys_prompt = system_override if system_override else SYSTEM_PROMPT
    # Добавляем факты из памяти пользователя
    if not system_override:
        facts = get_user_facts(user_id)
        if facts:
            memory_block = "\n\nЧто ты знаешь о пользователе (используй в диалоге):\n" + "\n".join(f"- {f}" for f in facts)
            sys_prompt = sys_prompt + memory_block
    messages = [{"role": "system", "content": sys_prompt}] + history

    if is_openrouter_model(selected_model):
        if not OPENROUTER_TOKEN:
            raise ValueError("OPENROUTER_TOKEN не задан")
        r = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_TOKEN}",
                "HTTP-Referer": "https://t.me/cosmoaibot",
                "X-Title": "Cosmo AI Bot",
            },
            json={"model": selected_model, "messages": messages},
            timeout=60
        )
    else:
        r = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {AI_TOKEN}"},
            json={"model": selected_model, "messages": messages},
            timeout=30
        )

    data = r.json()
    if "choices" not in data:
        raise ValueError(f"No choices in response: {data}")
    return data["choices"][0]["message"]["content"]

# ========= ОБРАБОТКА ЗАГРУЗКИ ПРОЕКТА =========
def extract_project_id_from_text(text: str) -> str | None:
    """Ищет 8-символьный ID проекта в тексте."""
    m = re.search(r'\b([A-F0-9]{8})\b', text.upper())
    if m:
        return m.group(1)
    return None

async def handle_project_load(update, context, project_id: str):
    row = load_project(project_id)
    if not row:
        await update.message.reply_text(
            f"<b>🏴‍☠️ Проект с ID <code>{project_id}</code> не найден.</b>",
            parse_mode="HTML"
        )
        return False

    _, uid, title, description, files_json, created_at = row
    files_data = json.loads(files_json)

    context.user_data["chat"] = True
    context.user_data["current_project_id"] = project_id

    files_list = "\n".join(f"  • {f['filename']}" for f in files_data)
    history_entry = f"Загружен проект '{title}' (ID: {project_id}). Файлы:\n{files_list}\n\nОписание: {description}"
    context.user_data["history"] = [{"role": "assistant", "content": history_entry}]

    await update.message.reply_text(
        f"<b>🏴‍☠️ Проект загружен!\n\n"
        f"📦 Название: {html.escape(title or 'Без названия')}\n"
        f"🆔 ID: <code>{project_id}</code>\n"
        f"📅 Создан: {created_at}\n\n"
        f"Файлы:\n{html.escape(files_list)}\n\n"
        f"Продолжаю работу с проектом. Что делать дальше?</b>",
        parse_mode="HTML",
        reply_markup=chat_keyboard()
    )
    return True

# ========= ЧАТ (текст) =========
async def chat(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    cursor.execute("SELECT banned FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    if not row or row[0] == 1:
        return

    text = update.message.text

    # ======= РЕЖИМ КЛОНИРОВАНИЯ: ждём текст для озвучки =======
    if context.user_data.get("voice_clone_step") == "waiting_text":
        sample_bytes = context.user_data.get("voice_clone_sample")
        if not sample_bytes:
            context.user_data.pop("voice_clone_step", None)
            await update.message.reply_text(
                "<b>🏴 Образец не найден. Начни сначала через меню.</b>",
                parse_mode="HTML", reply_markup=back()
            )
            return

        cursor.execute("SELECT requests FROM users WHERE user_id=?", (user_id,))
        r = cursor.fetchone()
        if not r or r[0] <= 0:
            await update.message.reply_text(
                "<b>🏴‍☠️ Запросы закончились! Купи ещё или пригласи друзей.</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🏴‍☠️ Купить", callback_data="buy")]])
            )
            return

        status_msg = await update.message.reply_text(
            "<b>🎓 Клонирую голос и генерирую речь...\n⏳ Это займёт 15–30 секунд.</b>",
            parse_mode="HTML"
        )
        try:
            audio_bytes = await clone_voice_and_speak(sample_bytes, text)

            cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
            conn.commit()

            await status_msg.delete()

            if audio_bytes:
                audio_buf = io.BytesIO(audio_bytes)
                audio_buf.name = "voice_clone.mp3"
                await update.message.reply_audio(
                    audio=audio_buf,
                    caption="<b>🎓 Голос склонирован! Вот результат.</b>",
                    parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup([
                        [InlineKeyboardButton("🎓 Ещё текст", callback_data="voice_gen")],
                        [InlineKeyboardButton("▪️ В меню", callback_data="menu")],
                    ])
                )
            else:
                await update.message.reply_text(
                    "<b>🏴 Не удалось сгенерировать голос.\n\n"
                    "Возможные причины:\n"
                    "• Образец слишком короткий (нужно 5+ сек)\n"
                    "• Проблемы с API ElevenLabs\n"
                    "• Превышен лимит бесплатного плана</b>",
                    parse_mode="HTML", reply_markup=back()
                )
        except Exception as e:
            logger.error(f"Voice clone text error: {e}")
            try:
                await status_msg.delete()
            except Exception:
                pass
            await update.message.reply_text(
                "<b>🏴 Ошибка клонирования голоса. Попробуй позже.</b>",
                parse_mode="HTML", reply_markup=back()
            )

        # Сбрасываем состояние после генерации
        context.user_data.pop("voice_clone_step", None)
        context.user_data.pop("voice_clone_sample", None)
        return

    # ======= ИНСТРУМЕНТЫ: обработка текстового ввода =======
    active_tool = context.user_data.get("tool")

    if active_tool == "weather":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>☁️ Получаю данные о погоде...</b>", parse_mode="HTML")
        result = get_weather(text.strip())
        await msg.delete()
        if result:
            await update.message.reply_text(result, parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("☁️ Другой город", callback_data="tool_weather")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await update.message.reply_text(
                "<b>☁️ Не удалось получить погоду. Проверь название города.</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "qr":
        context.user_data.pop("tool", None)
        qr_bytes = generate_qr(text.strip())
        if qr_bytes:
            buf = io.BytesIO(qr_bytes)
            buf.name = "qr.png"
            await update.message.reply_photo(photo=buf, caption=f"<b>🔗 QR-код для:</b>\n<code>{text[:100]}</code>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("🔗 Ещё QR", callback_data="tool_qr")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await update.message.reply_text("<b>🔗 Ошибка генерации QR-кода.</b>", parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "translate":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>🌐 Перевожу...</b>", parse_mode="HTML")
        prompt = f"Переведи следующий текст. Если в начале указан язык — переведи на него. Если не указан — переведи на английский. Верни ТОЛЬКО перевод без пояснений.\n\nТекст: {text}"
        translated = await ai_request_simple(prompt)
        await msg.delete()
        if translated:
            await update.message.reply_text(
                f"<b>🌐 Перевод:</b>\n\n{translated}",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("🌐 Ещё перевод", callback_data="tool_translate")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await update.message.reply_text("<b>🌐 Ошибка перевода.</b>", parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "search":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>🔍 Ищу в интернете...</b>", parse_mode="HTML")
        raw = search_web(text.strip())
        if raw:
            summary_prompt = f"На основе этих результатов поиска сделай краткую, информативную сводку на русском языке:\n\n{raw}"
            summary = await ai_request_simple(summary_prompt)
            await msg.delete()
            reply = f"<b>🔍 По запросу «{text[:50]}»:</b>\n\n{summary or raw[:1500]}"
            await update.message.reply_text(reply, parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("🔍 Новый поиск", callback_data="tool_search")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await msg.delete()
            await update.message.reply_text("<b>🔍 Ничего не найдено.</b>", parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "youtube":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>📺 Получаю субтитры видео...</b>", parse_mode="HTML")
        transcript = get_youtube_transcript(text.strip())
        if transcript:
            await msg.edit_text("<b>📺 Анализирую содержимое...</b>", parse_mode="HTML")
            summary_prompt = f"Сделай подробное краткое содержание этого видео на русском языке. Выдели главные мысли и ключевые моменты:\n\n{transcript}"
            summary = await ai_request_simple(summary_prompt)
            await msg.delete()
            await update.message.reply_text(
                f"<b>📺 Краткое содержание видео:</b>\n\n{summary or 'Не удалось обработать'}",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("📺 Другое видео", callback_data="tool_youtube")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await msg.delete()
            await update.message.reply_text(
                "<b>📺 Не удалось получить субтитры.\n\nВозможные причины:\n• Субтитры отключены\n• Видео недоступно\n• Неверная ссылка</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "pdf":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>📄 Генерирую PDF...</b>", parse_mode="HTML")
        content_prompt = f"Создай подробный структурированный документ по теме: {text}\n\nИспользуй заголовки (# Заголовок 1, ## Заголовок 2), параграфы. Пиши на русском языке."
        content = await ai_request_simple(content_prompt)
        if content:
            pdf_bytes = build_pdf(content, text[:50])
            await msg.delete()
            if pdf_bytes:
                buf = io.BytesIO(pdf_bytes)
                buf.name = "document.pdf"
                await update.message.reply_document(document=buf, filename="document.pdf",
                    caption=f"<b>📄 PDF: {text[:50]}</b>",
                    parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup([
                        [InlineKeyboardButton("📄 Ещё PDF", callback_data="tool_pdf")],
                        [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                    ]))
            else:
                await update.message.reply_text("<b>📄 Ошибка генерации PDF.</b>", parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        else:
            await msg.delete()
            await update.message.reply_text("<b>📄 Ошибка генерации контента.</b>", parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "excel":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>📊 Генерирую таблицу...</b>", parse_mode="HTML")
        table_prompt = f"Создай таблицу в формате Markdown (используй | для разделения колонок) по теме: {text}\n\nПервая строка — заголовки. Создай 10-15 строк с реалистичными данными."
        content = await ai_request_simple(table_prompt)
        if content:
            xl_bytes = build_excel(content, text[:30])
            await msg.delete()
            if xl_bytes:
                buf = io.BytesIO(xl_bytes)
                buf.name = "table.xlsx"
                await update.message.reply_document(document=buf, filename="table.xlsx",
                    caption=f"<b>📊 Excel: {text[:50]}</b>",
                    parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup([
                        [InlineKeyboardButton("📊 Ещё таблицу", callback_data="tool_excel")],
                        [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                    ]))
            else:
                await update.message.reply_text("<b>📊 Ошибка создания Excel.</b>", parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        else:
            await msg.delete()
            await update.message.reply_text("<b>📊 Ошибка генерации данных.</b>", parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "run_code":
        context.user_data.pop("tool", None)
        clean_code = text.strip().strip("```python").strip("```").strip()
        msg = await update.message.reply_text("<b>🐍 Запускаю код...</b>", parse_mode="HTML")
        result = await run_code_sandbox(clean_code)
        await msg.delete()
        await update.message.reply_text(
            f"<b>🐍 Результат:</b>\n\n{result}",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("🐍 Запустить ещё", callback_data="tool_code")],
                [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
            ]))
        return

    if active_tool == "memory":
        context.user_data.pop("tool", None)
        save_user_fact(user_id, text.strip())
        await update.message.reply_text(
            f"<b>🧠 Запомнил:</b>\n\n<i>{text[:200]}</i>\n\nAI будет использовать это в диалоге.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("🧠 Добавить ещё", callback_data="tool_memory")],
                [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
            ]))
        return

    if active_tool == "remind":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text("<b>⏰ Разбираю время напоминания...</b>", parse_mode="HTML")
        parse_prompt = (
            f"Сегодня: {__import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M')}.\n"
            f"Пользователь написал: «{text}»\n\n"
            "Извлеки дату/время и текст напоминания. Верни JSON:\n"
            '{"datetime": "YYYY-MM-DDTHH:MM", "text": "текст напоминания"}\n'
            "Только JSON, без пояснений."
        )
        parsed = await ai_request_simple(parse_prompt)
        await msg.delete()
        try:
            import json as _json
            parsed_clean = parsed.strip().strip("```json").strip("```").strip()
            data = _json.loads(parsed_clean)
            remind_dt = data["datetime"]
            remind_text = data["text"]
            add_reminder(user_id, remind_dt, remind_text)
            dt_display = remind_dt[:16].replace("T", " в ")
            await update.message.reply_text(
                f"<b>⏰ Напоминание установлено!</b>\n\n📅 Когда: <b>{dt_display}</b>\n📝 Что: {remind_text}",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("⏰ Ещё напоминание", callback_data="tool_remind")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        except Exception as e:
            logger.error(f"Reminder parse error: {e}, raw: {parsed}")
            await update.message.reply_text(
                "<b>⏰ Не смог разобрать время. Попробуй написать иначе:\n\n"
                "<i>«завтра в 15:00 позвонить»\n«через 2 часа сделать кофе»</i></b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    if active_tool == "music":
        context.user_data.pop("tool", None)
        msg = await update.message.reply_text(
            "<b>🎵 Генерирую музыку...\n⏳ Это займёт 30–60 секунд.</b>", parse_mode="HTML")
        audio_bytes = await generate_music_hf(text.strip())
        await msg.delete()
        if audio_bytes:
            buf = io.BytesIO(audio_bytes)
            buf.name = "music.wav"
            await update.message.reply_audio(audio=buf,
                caption=f"<b>🎵 Музыка по запросу:</b>\n<i>{text[:100]}</i>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("🎵 Ещё музыку", callback_data="tool_music")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await update.message.reply_text(
                "<b>🎵 Ошибка генерации. Возможно, модель загружается — попробуй через минуту.</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    # Проверка: просто ID проекта (или ID + ключевое слово)
    project_id_candidate = extract_project_id_from_text(text)
    if project_id_candidate:
        short_text = text.strip().upper()
        is_only_id = re.match(r'^[A-F0-9]{8}$', short_text)
        has_remember_kw = wants_remember(text)
        if is_only_id or has_remember_kw:
            await handle_project_load(update, context, project_id_candidate)
            return

    if not context.user_data.get("chat"):
        return

    cursor.execute("SELECT requests FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    req = row[0] if row else 0

    if req <= 0:
        await update.message.reply_text(
            "<b>🏴‍☠️ Запросы закончились! Купи ещё или пригласи друзей.</b>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("🏴‍☠️ Купить запросы", callback_data="buy")]
            ])
        )
        return

    # ========= КАРТИНКИ — обрабатываем ПЕРВЫМИ, до AI =========
    if wants_image(text):
        thinking_msg = await update.message.reply_text("<b>🎨 Генерирую картинку...</b>", parse_mode="HTML")
        try:
            image_bytes = await generate_image(text)
            await thinking_msg.delete()
            if image_bytes:
                cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
                conn.commit()
                buf = io.BytesIO(image_bytes)
                buf.name = "image.jpg"
                await update.message.reply_photo(
                    photo=buf,
                    caption=f"<b>🎨 Картинка по запросу: {html.escape(text[:100])}</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
            else:
                await update.message.reply_text(
                    "<b>🏴‍☠️ Не удалось сгенерировать картинку. Попробуй описать по-другому.</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
        except Exception as e:
            logger.error(f"Image error: {e}")
            try:
                await thinking_msg.delete()
            except Exception:
                pass
            await update.message.reply_text("<b>🏴‍☠️ Ошибка генерации картинки.</b>", parse_mode="HTML")
        return

    # ========= WORD ДОКУМЕНТ =========
    if wants_word(text):
        thinking_msg = await update.message.reply_text("<b>📝 Создаю Word документ...</b>", parse_mode="HTML")
        try:
            history = context.user_data.get("history", [])
            history.append({"role": "user", "content": text})
            ai_content = await ask_ai(user_id, history, system_override=WORD_SYSTEM_PROMPT)
            history.append({"role": "assistant", "content": ai_content})
            context.user_data["history"] = history[-20:]

            cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
            conn.commit()

            await thinking_msg.delete()
            doc_buf = build_word_doc(ai_content)
            if doc_buf:
                fname = "document.docx"
                await update.message.reply_document(
                    document=doc_buf,
                    filename=fname,
                    caption="<b>📝 Ваш Word документ готов!</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
            else:
                await update.message.reply_text(
                    f"<b>📝 Содержимое документа:\n\n{html.escape(ai_content[:3000])}</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
        except Exception as e:
            logger.error(f"Word error: {e}")
            try:
                await thinking_msg.delete()
            except Exception:
                pass
            await update.message.reply_text("<b>🏴‍☠️ Ошибка создания документа.</b>", parse_mode="HTML")
        return

    # ========= POWERPOINT ПРЕЗЕНТАЦИЯ =========
    if wants_ppt(text):
        thinking_msg = await update.message.reply_text("<b>📊 Создаю презентацию PowerPoint...</b>", parse_mode="HTML")
        try:
            history = context.user_data.get("history", [])
            history.append({"role": "user", "content": text})
            ai_content = await ask_ai(user_id, history, system_override=PPT_SYSTEM_PROMPT)
            history.append({"role": "assistant", "content": ai_content})
            context.user_data["history"] = history[-20:]

            cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
            conn.commit()

            await thinking_msg.delete()
            ppt_buf = build_ppt(ai_content, title=text[:50])
            if ppt_buf:
                await update.message.reply_document(
                    document=ppt_buf,
                    filename="presentation.pptx",
                    caption="<b>📊 Ваша презентация PowerPoint готова!</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
            else:
                await update.message.reply_text(
                    f"<b>📊 Содержимое презентации:\n\n{html.escape(ai_content[:3000])}</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
        except Exception as e:
            logger.error(f"PPT error: {e}")
            try:
                await thinking_msg.delete()
            except Exception:
                pass
            await update.message.reply_text("<b>🏴‍☠️ Ошибка создания презентации.</b>", parse_mode="HTML")
        return

    # ========= ОПРОС / ТЕСТ =========
    if wants_survey(text):
        thinking_msg = await update.message.reply_text("<b>📋 Создаю опрос...</b>", parse_mode="HTML")
        try:
            history = context.user_data.get("history", [])
            history.append({"role": "user", "content": text})
            ai_content = await ask_ai(user_id, history, system_override=SURVEY_SYSTEM_PROMPT)
            history.append({"role": "assistant", "content": ai_content})
            context.user_data["history"] = history[-20:]

            cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
            conn.commit()

            await thinking_msg.delete()
            survey_buf = build_survey_doc(ai_content)
            if survey_buf:
                await update.message.reply_document(
                    document=survey_buf,
                    filename="survey.docx",
                    caption="<b>📋 Ваш опрос готов!</b>",
                    parse_mode="HTML",
                    reply_markup=chat_keyboard()
                )
            else:
                for chunk in [ai_content[i:i+4000] for i in range(0, len(ai_content), 4000)]:
                    await update.message.reply_text(f"<b>{html.escape(chunk)}</b>", parse_mode="HTML")
                await update.message.reply_text("", reply_markup=chat_keyboard())
        except Exception as e:
            logger.error(f"Survey error: {e}")
            try:
                await thinking_msg.delete()
            except Exception:
                pass
            await update.message.reply_text("<b>🏴‍☠️ Ошибка создания опроса.</b>", parse_mode="HTML")
        return

    # ========= ОБЫЧНЫЙ AI ЧАТ / ZIP ПРОЕКТ =========
    history = context.user_data.get("history", [])
    history.append({"role": "user", "content": text})

    thinking_msg = await update.message.reply_text("<b>🏴‍☠️ Думаю...</b>", parse_mode="HTML")

    try:
        ai_reply = await ask_ai(user_id, history)
    except Exception as e:
        logger.error(f"AI error: {e}")
        await thinking_msg.delete()
        await update.message.reply_text("<b>🏴‍☠️ Ошибка при обращении к AI. Попробуй позже.</b>", parse_mode="HTML")
        return

    history.append({"role": "assistant", "content": ai_reply})
    context.user_data["history"] = history[-20:]

    cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
    conn.commit()

    await thinking_msg.delete()

    blocks = extract_code_blocks(ai_reply)
    if blocks and wants_zip(text):
        project_id = context.user_data.get("current_project_id") or generate_project_id()
        context.user_data["current_project_id"] = project_id

        project_type = detect_project_type(blocks)
        zip_name = get_project_zip_name(project_type, blocks)

        save_project(project_id, user_id, text[:80], ai_reply[:200], blocks)

        zip_buf = build_zip(blocks, project_id=project_id)
        names = ", ".join(f[0] for f in blocks)

        for chunk, _ in format_ai_reply_html(ai_reply):
            if len(chunk) > 4000:
                for i in range(0, len(chunk), 4000):
                    await update.message.reply_text(chunk[i:i+4000], parse_mode="HTML")
            else:
                await update.message.reply_text(chunk, parse_mode="HTML")

        await update.message.reply_document(
            document=zip_buf,
            filename=zip_name,
            caption=(
                f"<b>📦 Проект: {html.escape(names[:200])}\n\n"
                f"🆔 ID проекта: <code>{project_id}</code>\n"
                f"💡 Пришли этот ID чтобы продолжить работу с проектом в любое время!</b>"
            ),
            parse_mode="HTML",
            reply_markup=chat_keyboard()
        )
    else:
        parts = format_ai_reply_html(ai_reply)
        for idx, (chunk, is_code) in enumerate(parts):
            markup = chat_keyboard() if idx == len(parts) - 1 else None
            if len(chunk) > 4000:
                for i in range(0, len(chunk), 4000):
                    sub = chunk[i:i+4000]
                    await update.message.reply_text(sub, parse_mode="HTML", reply_markup=markup if i + 4000 >= len(chunk) else None)
            else:
                await update.message.reply_text(chunk, parse_mode="HTML", reply_markup=markup)

# ========= ГОЛОСОВЫЕ СООБЩЕНИЯ =========
async def voice_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    cursor.execute("SELECT banned FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    if not row or row[0] == 1:
        return

    # ======= РЕЖИМ КЛОНИРОВАНИЯ: ждём образец голоса =======
    if context.user_data.get("voice_clone_step") == "waiting_sample":
        status_msg = await update.message.reply_text(
            "<b>🎓 Образец получен, сохраняю...</b>", parse_mode="HTML"
        )
        try:
            voice = update.message.voice
            file = await context.bot.get_file(voice.file_id)
            with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as tmp:
                tmp_path = tmp.name
            await file.download_to_drive(tmp_path)
            with open(tmp_path, "rb") as f:
                sample_bytes = f.read()
            os.unlink(tmp_path)

            context.user_data["voice_clone_sample"] = sample_bytes
            context.user_data["voice_clone_step"] = "waiting_text"

            await status_msg.edit_text(
                "<b>🎓 Отлично! Образец голоса принят.\n\n"
                "Шаг 2 из 2: Теперь напиши текст, который нужно озвучить этим голосом.</b>",
                parse_mode="HTML"
            )
        except Exception as e:
            logger.error(f"Voice clone sample error: {e}")
            await status_msg.edit_text(
                "<b>🏴 Не удалось сохранить образец. Попробуй ещё раз.</b>",
                parse_mode="HTML"
            )
        return

    voice_only = context.user_data.get("voice_only", False)
    in_chat = context.user_data.get("chat", False)

    if not voice_only and not in_chat:
        await update.message.reply_text(
            "<b>🏴‍☠️ Нажми «Голосовые» для расшифровки или «Начать диалог» для чата с AI.</b>",
            parse_mode="HTML"
        )
        return

    thinking_msg = await update.message.reply_text("<b>🏴‍☠️ Расшифровываю...</b>", parse_mode="HTML")

    try:
        voice = update.message.voice
        file = await context.bot.get_file(voice.file_id)

        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as tmp:
            tmp_path = tmp.name

        await file.download_to_drive(tmp_path)

        with open(tmp_path, "rb") as audio_file:
            transcribe_r = requests.post(
                "https://api.groq.com/openai/v1/audio/transcriptions",
                headers={"Authorization": f"Bearer {AI_TOKEN}"},
                files={"file": ("voice.ogg", audio_file, "audio/ogg")},
                data={"model": "whisper-large-v3-turbo"},
                timeout=30
            )

        os.unlink(tmp_path)
        transcribe_data = transcribe_r.json()
        recognized_text = transcribe_data.get("text", "").strip()

        if not recognized_text:
            await thinking_msg.delete()
            await update.message.reply_text("<b>🏴‍☠️ Не смог распознать голос. Попробуй ещё раз.</b>", parse_mode="HTML")
            return

        if voice_only:
            await thinking_msg.delete()
            await update.message.reply_text(
                f"<b>🏴‍☠️ Расшифровка:\n\n{recognized_text}</b>",
                parse_mode="HTML",
                reply_markup=back()
            )
            return

        cursor.execute("SELECT requests FROM users WHERE user_id=?", (user_id,))
        row = cursor.fetchone()
        req = row[0] if row else 0

        if req <= 0:
            await thinking_msg.delete()
            await update.message.reply_text(
                "<b>🏴‍☠️ Запросы закончились! Купи ещё или пригласи друзей.</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("🏴‍☠️ Купить запросы", callback_data="buy")]
                ])
            )
            return

        await thinking_msg.edit_text(
            f"<b>🏴‍☠️ Распознано: {recognized_text}\n\nДумаю...</b>",
            parse_mode="HTML"
        )

        history = context.user_data.get("history", [])
        history.append({"role": "user", "content": recognized_text})

        ai_reply = await ask_ai(user_id, history)

        history.append({"role": "assistant", "content": ai_reply})
        context.user_data["history"] = history[-20:]

        cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
        conn.commit()

        await thinking_msg.delete()

        parts = format_ai_reply_html(ai_reply)
        for idx, (chunk, is_code) in enumerate(parts):
            markup = chat_keyboard() if idx == len(parts) - 1 else None
            if len(chunk) > 4000:
                for i in range(0, len(chunk), 4000):
                    sub = chunk[i:i+4000]
                    await update.message.reply_text(sub, parse_mode="HTML", reply_markup=markup if i + 4000 >= len(chunk) else None)
            else:
                await update.message.reply_text(chunk, parse_mode="HTML", reply_markup=markup)

    except Exception as e:
        logger.error(f"Voice error: {e}")
        try:
            await thinking_msg.delete()
        except Exception:
            pass
        await update.message.reply_text("<b>🏴‍☠️ Ошибка при обработке голосового. Попробуй позже.</b>", parse_mode="HTML")

# ========= ПРОМОКОДЫ =========
async def promo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    if not context.args:
        await update.message.reply_text("<b>🏴‍☠️ Введи: /promo КОД</b>", parse_mode="HTML")
        return

    code = context.args[0].upper()

    cursor.execute("SELECT requests, max_uses, uses_count FROM promo_codes WHERE code=?", (code,))
    promo_row = cursor.fetchone()

    if not promo_row:
        await update.message.reply_text("<b>🏴‍☠️ Промокод не найден.</b>", parse_mode="HTML")
        return

    req_amount, max_uses, uses_count = promo_row

    if uses_count >= max_uses:
        await update.message.reply_text("<b>🏴‍☠️ Промокод уже исчерпан.</b>", parse_mode="HTML")
        return

    cursor.execute("SELECT 1 FROM promo_uses WHERE code=? AND user_id=?", (code, user_id))
    if cursor.fetchone():
        await update.message.reply_text("<b>🏴‍☠️ Ты уже использовал этот промокод.</b>", parse_mode="HTML")
        return

    cursor.execute("UPDATE users SET requests = requests + ? WHERE user_id=?", (req_amount, user_id))
    cursor.execute("UPDATE promo_codes SET uses_count = uses_count + 1 WHERE code=?", (code,))
    cursor.execute("INSERT INTO promo_uses VALUES (?, ?)", (code, user_id))
    conn.commit()

    await update.message.reply_text(
        f"<b>🏴‍☠️ Промокод активирован!\n\n+{req_amount} запросов начислено.</b>",
        parse_mode="HTML"
    )

async def add_promo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    try:
        code = context.args[0].upper()
        req_amount = int(context.args[1])
        max_uses = int(context.args[2])

        cursor.execute(
            "INSERT OR REPLACE INTO promo_codes (code, requests, max_uses, uses_count) VALUES (?, ?, ?, 0)",
            (code, req_amount, max_uses)
        )
        conn.commit()

        await update.message.reply_text(
            f"<b>🏴‍☠️ Промокод создан!\n\n"
            f"Код: <code>{code}</code>\n"
            f"Запросов: {req_amount}\n"
            f"Макс. использований: {max_uses}</b>",
            parse_mode="HTML"
        )
    except Exception as e:
        await update.message.reply_text(
            f"<b>🏴‍☠️ Использование: /addpromo КОД КОЛИЧЕСТВО МАК_ИСПОЛЬЗОВАНИЙ\nОшибка: {e}</b>",
            parse_mode="HTML"
        )

# ========= ПОКУПКА ЛИМИТОВ =========
async def buy_limits(update: Update, context: ContextTypes.DEFAULT_TYPE, amount: int, usdt: float):
    user_id = update.effective_user.id

    try:
        r = requests.post(
            CRYPTO_API + "createInvoice",
            headers={"Crypto-Pay-API-Token": CRYPTO_TOKEN},
            json={
                "asset": "USDT",
                "amount": str(usdt),
                "description": f"Cosmo AI — {amount} запросов",
                "expires_in": 3600,
            },
            timeout=10
        )
        data = r.json()
        if not data.get("ok"):
            logger.error(f"Crypto Pay error: {data}")
            await update.message.reply_text("<b>🏴‍☠️ Ошибка создания счёта. Попробуй позже.</b>", parse_mode="HTML")
            return

        invoice = data["result"]
        invoice_id = str(invoice["invoice_id"])
        pay_url = invoice["bot_invoice_url"]

        cursor.execute("INSERT INTO payments VALUES (?, ?, ?, ?)", (invoice_id, user_id, amount, "pending"))
        conn.commit()

        await update.message.reply_text(
            f"<b>🏴‍☠️ Счёт создан!\n\n"
            f"🏴‍☠️ Сумма: {usdt} USDT\n"
            f"🏴‍☠️ Получишь: {amount} запросов\n\n"
            f"🏴‍☠️ После оплаты запросы будут начислены автоматически.</b>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("🏴‍☠️ Оплатить", url=pay_url)]
            ])
        )
    except Exception as e:
        logger.error(f"Payment error: {e}")
        await update.message.reply_text("<b>🏴‍☠️ Ошибка платёжной системы. Попробуй позже.</b>", parse_mode="HTML")

async def buy_100(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await buy_limits(update, context, 100, 1.0)

async def buy_300(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await buy_limits(update, context, 300, 2.5)

async def buy_1000(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await buy_limits(update, context, 1000, 7.0)

# ========= ADMIN КОМАНДЫ =========
async def cmd_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    try:
        username = context.args[0].replace("@", "")
        message = " ".join(context.args[1:])

        cursor.execute("SELECT user_id FROM users")
        users = cursor.fetchall()

        sent = 0
        for (uid,) in users:
            try:
                chat_info = await context.bot.get_chat(uid)
                if chat_info.username and chat_info.username.lower() == username.lower():
                    await context.bot.send_message(uid, f"<b>{message}</b>", parse_mode="HTML")
                    sent += 1
                    break
            except Exception:
                pass

        await update.message.reply_text(f"<b>🏴‍☠️ Отправлено {sent} пользователю(ям)</b>", parse_mode="HTML")
    except Exception as e:
        await update.message.reply_text(f"<b>🏴‍☠️ Ошибка: {e}</b>", parse_mode="HTML")

async def broadcast(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    message = " ".join(context.args)
    if not message:
        await update.message.reply_text("<b>🏴‍☠️ Укажи текст: /textTOP <сообщение></b>", parse_mode="HTML")
        return

    # Получаем всех пользователей ДО async операций — отдельным соединением
    tmp_conn = sqlite3.connect(_DB_PATH)
    tmp_cur = tmp_conn.cursor()
    tmp_cur.execute("SELECT user_id FROM users WHERE banned=0")
    all_users = [row[0] for row in tmp_cur.fetchall()]
    total = len(all_users)
    tmp_conn.close()

    status_msg = await update.message.reply_text(
        f"<b>🏴‍☠️ Начинаю рассылку {total} пользователям...</b>",
        parse_mode="HTML"
    )

    sent = 0
    failed = 0
    blocked = 0
    for uid in all_users:
        try:
            await context.bot.send_message(uid, f"<b>{message}</b>", parse_mode="HTML")
            sent += 1
        except Exception as e:
            err = str(e).lower()
            if "blocked" in err or "forbidden" in err or "deactivated" in err:
                blocked += 1
            else:
                failed += 1
        await asyncio.sleep(0.05)

    await status_msg.edit_text(
        f"<b>🏴‍☠️ Рассылка завершена!\n\n"
        f"🏴 Всего в базе: {total}\n"
        f"🎓 Доставлено: {sent}\n"
        f"🏴 Заблокировали бота: {blocked}\n"
        f"🏴 Другие ошибки: {failed}</b>",
        parse_mode="HTML"
    )

async def ban_user(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    try:
        username = context.args[0].replace("@", "")
        cursor.execute("SELECT user_id FROM users")
        users = cursor.fetchall()

        for (uid,) in users:
            if uid in ADMINS:
                continue
            try:
                chat_info = await context.bot.get_chat(uid)
                if chat_info.username and chat_info.username.lower() == username.lower():
                    cursor.execute("UPDATE users SET banned=1 WHERE user_id=?", (uid,))
                    conn.commit()
                    await update.message.reply_text(f"<b>🏴‍☠️ Пользователь @{username} заблокирован</b>", parse_mode="HTML")
                    return
            except Exception:
                pass

        await update.message.reply_text(f"<b>🏴‍☠️ Пользователь @{username} не найден</b>", parse_mode="HTML")
    except Exception as e:
        await update.message.reply_text(f"<b>🏴‍☠️ Ошибка: {e}</b>", parse_mode="HTML")

async def unban_user(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    try:
        username = context.args[0].replace("@", "")
        cursor.execute("SELECT user_id FROM users")
        users = cursor.fetchall()

        for (uid,) in users:
            try:
                chat_info = await context.bot.get_chat(uid)
                if chat_info.username and chat_info.username.lower() == username.lower():
                    cursor.execute("UPDATE users SET banned=0 WHERE user_id=?", (uid,))
                    conn.commit()
                    await update.message.reply_text(f"<b>🏴‍☠️ Пользователь @{username} разблокирован</b>", parse_mode="HTML")
                    return
            except Exception:
                pass

        await update.message.reply_text(f"<b>🏴‍☠️ Пользователь @{username} не найден</b>", parse_mode="HTML")
    except Exception as e:
        await update.message.reply_text(f"<b>🏴‍☠️ Ошибка: {e}</b>", parse_mode="HTML")

async def set_requests(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    try:
        amount = int(context.args[0])
        username = context.args[1].replace("@", "")

        cursor.execute("SELECT user_id FROM users")
        users = cursor.fetchall()

        for (uid,) in users:
            try:
                chat_info = await context.bot.get_chat(uid)
                if chat_info.username and chat_info.username.lower() == username.lower():
                    cursor.execute("UPDATE users SET requests = requests + ? WHERE user_id=?", (amount, uid))
                    conn.commit()
                    await update.message.reply_text(
                        f"<b>🏴‍☠️ Пользователю @{username} добавлено {amount} запросов</b>",
                        parse_mode="HTML"
                    )
                    return
            except Exception:
                pass

        await update.message.reply_text(f"<b>🏴‍☠️ Пользователь @{username} не найден</b>", parse_mode="HTML")
    except Exception as e:
        await update.message.reply_text(
            f"<b>🏴‍☠️ Использование: /set количество @username\nОшибка: {e}</b>",
            parse_mode="HTML"
        )

async def stats(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        return

    cursor.execute("SELECT COUNT(*) FROM users")
    total = cursor.fetchone()[0]

    cursor.execute("SELECT COUNT(*) FROM users WHERE banned=1")
    banned = cursor.fetchone()[0]

    cursor.execute("SELECT COUNT(*) FROM payments WHERE status='paid'")
    paid = cursor.fetchone()[0]

    cursor.execute("SELECT SUM(total_used) FROM users")
    all_used = cursor.fetchone()[0] or 0

    cursor.execute("SELECT COUNT(*) FROM promo_codes")
    promo_count = cursor.fetchone()[0]

    cursor.execute("SELECT SUM(uses_count) FROM promo_codes")
    promo_uses = cursor.fetchone()[0] or 0

    cursor.execute("SELECT COUNT(*) FROM projects")
    proj_count = cursor.fetchone()[0]

    await update.message.reply_text(
        f"<b>🏴‍☠️ Статистика бота\n\n"
        f"🏴‍☠️ Всего пользователей: {total}\n"
        f"🏴‍☠️ Заблокировано: {banned}\n"
        f"🏴‍☠️ Оплаченных счетов: {paid}\n"
        f"🏴‍☠️ Всего запросов к AI: {all_used}\n"
        f"🏴‍☠️ Промокодов создано: {promo_count}\n"
        f"🏴‍☠️ Промокодов активировано: {promo_uses}\n"
        f"🏴‍☠️ Проектов сохранено: {proj_count}</b>",
        parse_mode="HTML"
    )

# ========= ПРОВЕРКА ОПЛАТЫ =========
async def check_payments(context: ContextTypes.DEFAULT_TYPE):
    try:
        r = requests.get(
            CRYPTO_API + "getInvoices",
            headers={"Crypto-Pay-API-Token": CRYPTO_TOKEN},
            timeout=10
        )
        data = r.json()

        if not data.get("ok"):
            return

        for inv in data["result"].get("items", []):
            if inv["status"] == "paid":
                invoice_id = str(inv["invoice_id"])

                cursor.execute("SELECT * FROM payments WHERE invoice_id=?", (invoice_id,))
                p = cursor.fetchone()

                if p and p[3] == "pending":
                    user_id = p[1]
                    amount = p[2]
                    cursor.execute("UPDATE users SET requests = requests + ? WHERE user_id=?", (amount, user_id))
                    cursor.execute("UPDATE payments SET status='paid' WHERE invoice_id=?", (invoice_id,))
                    conn.commit()

                    try:
                        await context.bot.send_message(
                            user_id,
                            f"<b>🏴‍☠️ Оплата получена! Вам начислено {amount} запросов.</b>",
                            parse_mode="HTML"
                        )
                    except Exception:
                        pass
    except Exception as e:
        logger.error(f"Payment check error: {e}")

# ========= ОБРАБОТЧИК ФОТОГРАФИЙ =========
async def photo_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    cursor.execute("SELECT banned FROM users WHERE user_id=?", (user_id,))
    row = cursor.fetchone()
    if not row or row[0] == 1:
        return

    active_tool = context.user_data.get("tool")
    caption = (update.message.caption or "").strip()

    # Скачиваем фото
    try:
        photo = update.message.photo[-1]
        file = await context.bot.get_file(photo.file_id)
        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
            tmp_path = tmp.name
        await file.download_to_drive(tmp_path)
        with open(tmp_path, "rb") as f:
            image_bytes = f.read()
        os.unlink(tmp_path)
    except Exception as e:
        logger.error(f"Photo download error: {e}")
        await update.message.reply_text("<b>🏴 Не удалось загрузить фото.</b>", parse_mode="HTML")
        return

    # Режим редактирования фото
    if active_tool == "photo_edit":
        if not caption:
            await update.message.reply_text(
                "<b>🖌 Добавь подпись к фото с командой редактирования.\n\n"
                "Пример: отправь фото с подписью «черно-белое»</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
            return
        edited = edit_photo_pillow(image_bytes, caption)
        if edited:
            buf = io.BytesIO(edited)
            buf.name = "edited.jpg"
            await update.message.reply_photo(photo=buf,
                caption=f"<b>🖌 Готово: {caption}</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("🖌 Ещё правку", callback_data="tool_edit")],
                    [InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]
                ]))
        else:
            await update.message.reply_text(
                "<b>🖌 Команда не распознана. Попробуй:\nчерно-белое, зеркало, ярче, темнее, контраст, размытие, резкость, инвертировать</b>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("▪️ Инструменты", callback_data="tools")]]))
        return

    # Режим анализа фото (или авто при любом фото)
    cursor.execute("SELECT requests FROM users WHERE user_id=?", (user_id,))
    req_row = cursor.fetchone()
    if not req_row or req_row[0] <= 0:
        await update.message.reply_text(
            "<b>🏴‍☠️ Запросы закончились! Купи ещё или пригласи друзей.</b>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🏴‍☠️ Купить", callback_data="buy")]]))
        return

    analyze_prompt = caption if caption else "Опиши подробно что изображено на этом фото"
    msg = await update.message.reply_text("<b>🖼 Анализирую изображение...</b>", parse_mode="HTML")

    result = await analyze_image_with_ai(image_bytes, analyze_prompt)
    cursor.execute("UPDATE users SET requests = requests - 1, total_used = total_used + 1 WHERE user_id=?", (user_id,))
    conn.commit()

    await msg.delete()
    if result:
        await update.message.reply_text(
            f"<b>🖼 Анализ изображения:</b>\n\n{result}",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("🖼 Анализ ещё", callback_data="tool_img")],
                [InlineKeyboardButton("▪️ В меню", callback_data="menu")]
            ]))
    else:
        await update.message.reply_text(
            "<b>🖼 Не удалось проанализировать изображение. Попробуй ещё раз.</b>",
            parse_mode="HTML", reply_markup=back())

    if active_tool == "img_analyze":
        context.user_data.pop("tool", None)

# ========= ЗАПУСК =========
async def post_init(application):
    await application.bot.set_my_commands([
        ("start",    "Главное меню"),
        ("promo",    "Активировать промокод"),
        ("buy_100",  "Купить 100 запросов — 1 USDT"),
        ("buy_300",  "Купить 300 запросов — 2.5 USDT"),
        ("buy_1000", "Купить 1000 запросов — 7 USDT"),
    ])

def main():
    app = ApplicationBuilder().token(BOT_TOKEN).post_init(post_init).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("promo", promo))
    app.add_handler(CommandHandler("buy_100", buy_100))
    app.add_handler(CommandHandler("buy_300", buy_300))
    app.add_handler(CommandHandler("buy_1000", buy_1000))
    app.add_handler(CommandHandler("text", cmd_text))
    app.add_handler(CommandHandler("textTOP", broadcast))
    app.add_handler(CommandHandler("ban", ban_user))
    app.add_handler(CommandHandler("unban", unban_user))
    app.add_handler(CommandHandler("set", set_requests))
    app.add_handler(CommandHandler("stats", stats))
    app.add_handler(CommandHandler("addpromo", add_promo))

    app.add_handler(CallbackQueryHandler(buttons))
    app.add_handler(MessageHandler(filters.VOICE, voice_handler))
    app.add_handler(MessageHandler(filters.PHOTO, photo_handler))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, chat), group=1)

    app.job_queue.run_repeating(check_payments, interval=15)
    app.job_queue.run_repeating(check_reminders_job, interval=30)

    logger.info("🏴‍☠️ Cosmo AI Bot запущен!")
    app.run_polling(drop_pending_updates=True)

if __name__ == "__main__":
    main()
